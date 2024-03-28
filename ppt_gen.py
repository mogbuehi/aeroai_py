from pptx import Presentation
from pptx.util import Inches

import streamlit as st
from pptx import Presentation
from pptx.util import Inches

from pptx import Presentation
from pptx.util import Inches

from openai import OpenAI
import os
from dotenv import load_dotenv, find_dotenv

import json

def text_2_json (ppt_text):
    load_dotenv(find_dotenv()) 
    api_key = os.getenv('OPENAI_API_KEY')
    client = OpenAI(api_key=api_key)

    response = client.chat.completions.create(
    model="gpt-3.5-turbo-1106",
    response_format={ "type": "json_object" },
    messages=[
        {"role": "system", "content": """You are a helpful assistant designed to output JSON for a powerpoint slide. ONLY return the JSON, nothing else
        The schema to be returned is the following: 
        {\'title':'Put the title of the slide here', 'paragraph': ['paragraph 1','paragraph 2',...] } 
        
        example: 
        Title: Vehicles
        Paragraph 1: I like to fly planes
        Paragraph 2: I like to drive cars
        Paragraph 3: I like to drive boats
        
        
        JSON output:  {\'title':'Vehicles', 'paragraph': ['I like to fly planes','I like to drive cars','I like to drive boats'] } """},
        
        {"role": "user", "content": ppt_text}
        ]
    )
    json_str = response.choices[0].message.content
    print(json_str) #debugging
    json_data = json.loads(json_str)
    return json_data
    
def ppt_gen(ppt_text, ppt_path='presentation.pptx' ,image_path=None):      
    prs = Presentation(ppt_path)
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    # take ppt_text and turn it into json
    json_data = text_2_json(ppt_text)

    # parse json for the title of the slide
    title_shape.text = json_data['title']
    
    paragraph_list = json_data['paragraph']
    

    # by default first level is not defined (essentially p.level = 0)
    tf = body_shape.text_frame
    for paragraph in (paragraph_list):
        p = tf.add_paragraph() 
        p.text = paragraph

    # p = tf.add_paragraph()
    # p.text = 'Use _TextFrame.text for first bullet'
    # p.level = 1

    # p = tf.add_paragraph()
    # p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
    # p.level = 2

    # # Save the PowerPoint presentation with a dynamic filename
    # if image_path:
    #     presentation_filename = "presentation_with_image.pptx"
    # else:
    #     presentation_filename = "presentation_without_image.pptx"
    # # actually this should just add slides to a presentation
    prs.save(ppt_path) 

    # Example usage:
    # ppt_gen("Custom text goes here", "path_to_image.png")
    return ppt_text




from io import BytesIO
import tempfile
import os

def convert_pptx_to_images(pptx_path):
    prs = Presentation(pptx_path)
    images = []
    for i, slide in enumerate(prs.slides):
        # Use a BytesIO buffer to save the image
        image_stream = BytesIO()
        # Use a temporary file to save the image
        with tempfile.NamedTemporaryFile(delete=True, suffix='.png') as tmpfile:
            slide.shapes.save_slide(slide, tmpfile.name)
            image_stream.write(tmpfile.read())
        images.append(image_stream)
    return images


if __name__ == "__main__":
    ppt_text = """
Title: "The Possibilities"

- Uncharted territory
- Boundless creativity
- Limitless potential
- A world of opportunities
- The road less traveled
"""
    ppt_gen(ppt_text=ppt_text)





    # may need a format parser for input going into the function, or use f string to construct input going into the function
